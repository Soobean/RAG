"""
문서 처리 모듈: GPT-4o를 활용한 문서 파일 처리 및 검색 가능한 형태로 변환
"""
import os
import logging
import json
import io
import base64
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
from PIL import Image

import fitz  # PyMuPDF

# 로깅 설정
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(name)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """GPT-4o 기반 문서 처리 및 인덱싱 클래스"""

    def __init__(self, config: Dict[str, Any], search_engine):
        """
        문서 처리기 초기화

        Args:
            config: 애플리케이션 설정
            search_engine: 검색 엔진 인스턴스
        """
        self.config = config
        self.search_engine = search_engine
        self.openai_client = search_engine.openai_client

        self.processing_config = config.get('processing', {})
        self.max_image_width = self.processing_config.get('max_image_width', 1000)
        self.image_quality = self.processing_config.get('default_compression_quality', 85)
        self.max_image_size_kb = self.processing_config.get('max_image_size_kb', 1024)

        self.multimodal_model = config.get('openai', {}).get('multimodal_model', 'gpt-4o')

        self.current_document_name = ""

    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        GPT-4o 기반 문서 처리 메인 함수

        Args:
            file_path: 처리할 문서 경로

        Returns:
            처리 결과 정보
        """
        try:
            self.current_document_name = os.path.basename(file_path).split('.')[0]
            file_extension = os.path.splitext(file_path)[1].lower()

            logger.info(f"문서 처리 시작: {self.current_document_name}{file_extension}")

            if file_extension == '.pdf':
                processed_pages = self._process_pdf(file_path)
            elif file_extension == '.pptx':
                processed_pages = self._process_pptx(file_path)
            else:
                raise ValueError(f"지원되지 않는 파일 형식: {file_extension}")

            if hasattr(self.search_engine, 'migrate_to_document_structure_in_place'):
                try:
                    self.search_engine.migrate_to_document_structure_in_place()
                except Exception as e:
                    logger.warning(f"문서 구조 마이그레이션 중 오류: {e}")

            return {
                "document_name": self.current_document_name,
                "pages_processed": len(processed_pages),
                "pages": processed_pages,
                "status": "success"
            }

        except Exception as e:
            logger.error(f"문서 처리 중 오류 발생: {e}", exc_info=True)
            return {
                "document_name": self.current_document_name if hasattr(self, 'current_document_name') else "unknown",
                "pages_processed": 0,
                "status": "error",
                "error": str(e)
            }

    def _process_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        PDF 파일 처리

        Args:
            pdf_path: PDF 파일 경로

        Returns:
            처리된 페이지 목록
        """
        doc = fitz.open(pdf_path)
        processed_pages = []

        total_pages = len(doc)
        logger.info(f"PDF 처리 시작: {pdf_path}, 총 {total_pages}페이지")

        for page_idx in range(total_pages):
            page_num = page_idx + 1
            logger.info(f"페이지 {page_num}/{total_pages} 처리 중...")

            try:
                page = doc[page_idx]

                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                page_image = self._process_and_encode_image(img)

                analysis_result = self._analyze_with_gpt4o(page_image)

                processed_content = self._process_page_content(page, analysis_result, page_num)

                processed_content['combined_embedding'] = self._create_optimized_embedding(processed_content)

                document = self._prepare_for_storage(processed_content)
                self.search_engine.upsert_document(document)

                processed_pages.append({
                    'page_number': page_num,
                    'summary': processed_content.get('page_summary', '')
                })

                logger.info(f"페이지 {page_num} 처리 완료")

            except Exception as e:
                logger.error(f"페이지 {page_num} 처리 중 오류: {e}", exc_info=True)

        return processed_pages

    def _process_pptx(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        PPTX 파일 처리

        Args:
            pptx_path: PPTX 파일 경로

        Returns:
            처리된 페이지 목록
        """
        from pptx import Presentation

        prs = Presentation(pptx_path)
        processed_pages = []

        total_slides = len(prs.slides)
        logger.info(f"PPTX 처리 시작: {pptx_path}, 총 {total_slides}슬라이드")

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            logger.info(f"슬라이드 {slide_num}/{total_slides} 처리 중...")

            try:
                slide_image = self._extract_slide_as_image(slide, pptx_path, slide_idx)

                analysis_result = self._analyze_with_gpt4o(slide_image)

                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"

                processed_content = {
                    'folder_name': self.current_document_name,
                    'page_number': str(slide_num),
                    'text_content': slide_text,
                    'page_summary': analysis_result.get('page_summary', ''),
                    'elements': analysis_result.get('elements', []),
                    'images': [],
                }

                for element in analysis_result.get('elements', []):
                    if element.get('type') == 'image':
                        image_info = {
                            'image': slide_image,  # 전체 슬라이드 이미지 사용
                            'description': element.get('description', ''),
                            'related_text_ids': element.get('related_text_ids', [])
                        }
                        processed_content['images'].append(image_info)

                if not processed_content['images']:
                    processed_content['images'].append({
                        'image': slide_image,
                        'description': '슬라이드 전체 이미지',
                        'related_text_ids': []
                    })

                processed_content['combined_embedding'] = self._create_optimized_embedding(processed_content)

                document = self._prepare_for_storage(processed_content)
                self.search_engine.upsert_document(document)

                processed_pages.append({
                    'page_number': slide_num,
                    'summary': processed_content.get('page_summary', '')
                })

                logger.info(f"슬라이드 {slide_num} 처리 완료")

            except Exception as e:
                logger.error(f"슬라이드 {slide_num} 처리 중 오류: {e}", exc_info=True)

        return processed_pages

    def _extract_slide_as_image(self, slide, pptx_path, slide_idx):
        """
        PPTX 슬라이드를 이미지로 변환

        Args:
            slide: pptx 슬라이드 객체
            pptx_path: PPTX 파일 경로
            slide_idx: 슬라이드 인덱스

        Returns:
            Base64 인코딩된 이미지
        """
        try:
            try:
                from pptx_export import export_slides

                temp_dir = os.path.join(os.getcwd(), "temp_slides")
                os.makedirs(temp_dir, exist_ok=True)

                export_slides(pptx_path, temp_dir, slide_numbers=[slide_idx + 1])

                image_path = os.path.join(temp_dir, f"slide_{slide_idx + 1}.png")
                with open(image_path, 'rb') as f:
                    img_data = f.read()

                img = Image.open(io.BytesIO(img_data))
                return self._process_and_encode_image(img)
            except ImportError:
                pass

            try:
                import subprocess
                import glob

                temp_dir = os.path.join(os.getcwd(), "temp_slides")
                os.makedirs(temp_dir, exist_ok=True)
                output_path = os.path.join(temp_dir, "slide.pdf")

                subprocess.call([
                    'soffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    pptx_path
                ])

                pdf_files = glob.glob(os.path.join(temp_dir, "*.pdf"))
                if pdf_files:
                    doc = fitz.open(pdf_files[0])
                    if slide_idx < len(doc):
                        page = doc[slide_idx]
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                        doc.close()
                        return self._process_and_encode_image(img)
            except Exception:
                pass

            from PIL import ImageDraw, ImageFont
            img = Image.new('RGB', (1024, 768), color='white')
            draw = ImageDraw.Draw(img)

            font = None
            try:
                font = ImageFont.load_default()
            except:
                pass

            draw.text((20, 20), f"슬라이드 {slide_idx + 1}", fill='black', font=font)

            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            y_position = 60
            for line in text.split('\n'):
                draw.text((20, y_position), line, fill='black', font=font)
                y_position += 20
                if y_position > 700:  # 이미지 높이 제한
                    break

            return self._process_and_encode_image(img)

        except Exception as e:
            logger.error(f"슬라이드 이미지 추출 오류: {e}", exc_info=True)
            img = Image.new('RGB', (800, 600), color='white')
            draw = ImageDraw.Draw(img)
            draw.text((300, 300), f"슬라이드 {slide_idx + 1} - 이미지 추출 실패", fill='black')
            return self._process_and_encode_image(img)

    def _analyze_with_gpt4o(self, image_data):
        """
        GPT-4o를 사용하여 이미지 분석

        Args:
            image_data: Base64 인코딩된 이미지

        Returns:
            분석 결과
        """
        if not self.openai_client:
            logger.error("OpenAI 클라이언트가 초기화되지 않았습니다.")
            return {"elements": [], "page_summary": "분석 실패: OpenAI 클라이언트 없음"}

        if not image_data or not image_data.startswith('data:image'):
            logger.error("유효하지 않은 이미지 데이터")
            return {"elements": [], "page_summary": "분석 실패: 이미지 데이터 없음"}

        try:
            messages = [
                {
                    "role": "system",
                    "content": """
                    문서 분석 전문가로서 페이지 이미지를 분석하세요. 각 요소(이미지, 텍스트)를 식별하고 
                    다음 정보를 제공하세요:

                    1. 각 이미지 요소의 대략적 위치, 유형(스크린샷, 도표, 로고 등), 내용 설명
                    2. 각 텍스트 요소의 대략적 위치, 유형(제목, 본문, 캡션 등), 내용 요약
                    3. 텍스트-이미지 간 연관 관계(예: "이 텍스트는 왼쪽 도표를 설명함")

                    JSON 형식으로 응답하세요:
                    {
                      "elements": [
                        {
                          "type": "image",
                          "coordinates": {"x1": 0, "y1": 0, "x2": 100, "y2": 100},
                          "content_type": "screenshot|chart|diagram|photo|logo",
                          "description": "이미지에 대한 간결한 설명",
                          "related_text_ids": [1, 3]
                        },
                        {
                          "type": "text",
                          "coordinates": {"x1": 0, "y1": 0, "x2": 100, "y2": 100},
                          "content_type": "title|heading|body|caption|list",
                          "summary": "텍스트 내용 요약"
                        }
                      ],
                      "page_summary": "페이지 전체 내용 요약"
                    }
                    """
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "이 문서 페이지를 분석하여 이미지와 텍스트 요소를 구분하고 상세 정보를 제공해주세요."
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": image_data}
                        }
                    ]
                }
            ]

            response = self.openai_client.chat.completions.create(
                model=self.multimodal_model,
                messages=messages,
                response_format={"type": "json_object"},
                temperature=0.2
            )

            result = json.loads(response.choices[0].message.content)
            logger.info(f"GPT-4o 분석 결과: {len(result.get('elements', []))}개 요소 감지")
            return result

        except Exception as e:
            logger.error(f"GPT-4o 분석 오류: {e}", exc_info=True)
            return {"elements": [], "page_summary": f"분석 실패: {str(e)}"}

    def _process_page_content(self, page, analysis_result, page_number):
        """
        GPT-4o 분석 결과를 바탕으로 페이지 콘텐츠 처리

        Args:
            page: 페이지 객체
            analysis_result: GPT-4o 분석 결과
            page_number: 페이지 번호

        Returns:
            처리된 페이지 콘텐츠
        """
        width, height = page.rect.width, page.rect.height

        processed_content = {
            'folder_name': self.current_document_name,
            'page_number': str(page_number),
            'page_summary': analysis_result.get('page_summary', ''),
            'text_content': '',
            'images': [],
            'elements': []
        }

        for idx, element in enumerate(analysis_result.get('elements', [])):
            element_type = element.get('type')
            coords = element.get('coordinates', {})

            rect = fitz.Rect(
                coords.get('x1', 0) * width / 100,
                coords.get('y1', 0) * height / 100,
                coords.get('x2', 100) * width / 100,
                coords.get('y2', 100) * height / 100
            )

            element_data = {
                'id': idx,
                'type': element_type,
                'content_type': element.get('content_type', ''),
                'coordinates': coords
            }

            if element_type == 'image':
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=rect)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    image_data = self._process_and_encode_image(img)

                    image_info = {
                        'image': image_data,
                        'description': element.get('description', ''),
                        'related_text_ids': element.get('related_text_ids', [])
                    }

                    processed_content['images'].append(image_info)
                    element_data['description'] = element.get('description', '')
                    element_data['image_index'] = len(processed_content['images']) - 1
                except Exception as e:
                    logger.warning(f"이미지 요소 추출 오류: {e}")

            elif element_type == 'text':
                try:
                    text_content = page.get_text("text", clip=rect)

                    if not text_content.strip():
                        continue

                    element_data['content'] = text_content
                    element_data['summary'] = element.get('summary', '')
                    element_data['related_image_ids'] = element.get('related_image_ids', [])

                    processed_content['text_content'] += text_content + "\n\n"
                except Exception as e:
                    logger.warning(f"텍스트 요소 추출 오류: {e}")

            processed_content['elements'].append(element_data)

        if not processed_content['images']:
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                page_image = self._process_and_encode_image(img)

                processed_content['images'].append({
                    'image': page_image,
                    'description': '페이지 전체 이미지',
                    'related_text_ids': []
                })
            except Exception as e:
                logger.warning(f"전체 페이지 이미지 추출 오류: {e}")

        if not processed_content['text_content'].strip():
            processed_content['text_content'] = "텍스트 내용 없음"

        return processed_content

    def _process_and_encode_image(self, img):
        """
        이미지 처리 및 Base64 인코딩

        Args:
            img: PIL 이미지 객체

        Returns:
            Base64 인코딩된 이미지 데이터
        """
        try:
            if img.width > self.max_image_width:
                ratio = self.max_image_width / img.width
                new_height = int(img.height * ratio)
                img = img.resize((self.max_image_width, new_height), Image.LANCZOS)

            buffer = io.BytesIO()

            quality = self.image_quality
            img.save(buffer, format="JPEG", quality=quality, optimize=True)

            max_size_bytes = self.max_image_size_kb * 1024
            if buffer.tell() > max_size_bytes:
                for lower_quality in [70, 50, 30]:
                    buffer.seek(0)
                    buffer.truncate(0)
                    img.save(buffer, format="JPEG", quality=lower_quality, optimize=True)
                    if buffer.tell() <= max_size_bytes:
                        break

            buffer.seek(0)
            encoded_image = base64.b64encode(buffer.read()).decode('utf-8')
            return f"data:image/jpeg;base64,{encoded_image}"

        except Exception as e:
            logger.error(f"이미지 처리 오류: {e}", exc_info=True)
            return ""

    def _create_optimized_embedding(self, processed_content):
        """
        최적화된 임베딩 텍스트 생성

        Args:
            processed_content: 처리된 페이지 콘텐츠

        Returns:
            임베딩 벡터
        """
        embedding_text = f"""
        문서: {processed_content.get('folder_name', '')}
        페이지: {processed_content.get('page_number', '')}
        요약: {processed_content.get('page_summary', '')}

        텍스트 내용:
        {processed_content.get('text_content', '')[:4000]}

        이미지 설명:
        """

        for img in processed_content.get('images', []):
            embedding_text += f"- {img.get('description', '')}\n"

        return self.search_engine.generate_embedding(embedding_text)

    def _prepare_for_storage(self, processed_content):
        """
        저장용 문서 형식 준비

        Args:
            processed_content: 처리된 페이지 콘텐츠

        Returns:
            저장용 문서 객체
        """

        document = {
            '_id': f"{processed_content['folder_name']}_page_{processed_content['page_number']}",
            'folder_name': processed_content['folder_name'],
            'page_number': processed_content['page_number'],
            'description': processed_content['text_content'],
            'page_summary': processed_content['page_summary'],
            'images': processed_content['images'],
            'elements': processed_content.get('elements', []),
            'combined_embedding': processed_content['combined_embedding'],
            'created_at': datetime.now(),
            'content_type': 'processed_page'
        }

        return document