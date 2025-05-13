import os
import tempfile
import subprocess
import logging
import io
import base64
from typing import Dict, Any, List
from PIL import Image, ImageDraw, ImageFont

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTXProcessor:
    """PPTX 파일 처리 클래스"""

    def __init__(self, config):
        """
        PPTX 처리기 초기화
        """
        self.config = config
        processing_config = config.get('processing', {})
        self.max_image_width = processing_config.get('max_image_width', 1000)
        self.image_quality = processing_config.get('default_compression_quality', 85)

    def process(self, pptx_path, folder_name):
        """
        PPTX 처리 메인 함수
        """
        try:
            from pptx import Presentation

            prs = Presentation(pptx_path)
            pages_data = []

            total_slides = len(prs.slides)
            logger.info(f"PPTX 처리 시작: {pptx_path}, 총 {total_slides}슬라이드")

            for slide_idx, slide in enumerate(prs.slides):
                slide_num = slide_idx + 1
                logger.info(f"슬라이드 {slide_num}/{total_slides} 처리 중...")

                try:
                    slide_text = self._extract_slide_text(slide)

                    slide_image = self._extract_slide_as_image(slide, pptx_path, slide_idx)

                    slide_data = {
                        "folder_name": folder_name,
                        "page_number": str(slide_num),
                        "text_content": slide_text,
                        "page_summary": f"슬라이드 {slide_num}",
                        "images": [{
                            'image': slide_image,
                            'description': '슬라이드 전체 이미지',
                            'related_text_ids': []
                        }],
                        "elements": [
                            {
                                "id": 0,
                                "type": "text",
                                "content_type": "body",
                                "content": slide_text,
                                "summary": "슬라이드 텍스트"
                            }
                        ]
                    }

                    pages_data.append(slide_data)
                    logger.info(f"슬라이드 {slide_num} 처리 완료")

                except Exception as e:
                    logger.error(f"슬라이드 {slide_num} 처리 중 오류: {e}")
                    pages_data.append({
                        "folder_name": folder_name,
                        "page_number": str(slide_num),
                        "text_content": f"슬라이드 {slide_num} 처리 오류",
                        "page_summary": f"슬라이드 {slide_num} - 처리 실패",
                        "images": [],
                        "elements": []
                    })

            return pages_data

        except ImportError as e:
            logger.error(f"python-pptx 임포트 오류: {e}")
            return self._process_pptx_alternative(pptx_path, folder_name)

        except Exception as e:
            logger.error(f"PPTX 처리 중 오류: {e}")
            return [{
                "folder_name": folder_name,
                "page_number": "1",
                "text_content": f"PPTX 처리 오류: {str(e)}",
                "page_summary": "PPTX 처리 실패",
                "images": [],
                "elements": []
            }]

    def _extract_slide_text(self, slide):
        """
        슬라이드에서 텍스트 추출
        """
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"

        return slide_text.strip()

    def _extract_slide_as_image(self, slide, pptx_path, slide_idx):
        """
        슬라이드를 이미지로 변환
        """
        try:
            return self._extract_via_libreoffice(pptx_path, slide_idx)
        except Exception as e:
            logger.warning(f"LibreOffice 변환 실패: {e}")
            try:
                return self._create_text_image(slide, slide_idx)
            except Exception as e:
                logger.warning(f"텍스트 이미지 생성 실패: {e}")
                return self._create_default_image(slide_idx)

    def _extract_via_libreoffice(self, pptx_path, slide_idx):
        """
        LibreOffice를 사용해 슬라이드 추출
        """
        temp_dir = tempfile.mkdtemp()
        output_pdf = os.path.join(temp_dir, "slides.pdf")

        subprocess.call([
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            pptx_path
        ])

        try:
            import fitz  # PyMuPDF
            doc = fitz.open(output_pdf)

            if slide_idx < len(doc):
                page = doc[slide_idx]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                doc.close()

                import shutil
                shutil.rmtree(temp_dir)

                return self._process_and_encode_image(img)

            raise Exception("슬라이드를 찾을 수 없습니다")
        except Exception as e:
            import shutil
            shutil.rmtree(temp_dir)
            raise e

    def _create_text_image(self, slide, slide_idx):
        """
        슬라이드 텍스트 기반 이미지 생성
        """
        width, height = 1024, 768
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        font = ImageFont.load_default()

        draw.text((20, 20), f"슬라이드 {slide_idx + 1}", fill='black', font=font)

        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        y_position = 60
        for line in text.split('\n'):
            draw.text((20, y_position), line, fill='black', font=font)
            y_position += 20
            if y_position > height - 40:
                break

        return self._process_and_encode_image(img)

    def _create_default_image(self, slide_idx):
        """
        기본 슬라이드 이미지 생성
        """
        width, height = 1024, 768
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        font = ImageFont.load_default()

        draw.text((width // 2 - 100, height // 2 - 20), f"슬라이드 {slide_idx + 1}", fill='black', font=font)
        draw.text((width // 2 - 150, height // 2 + 20), "이미지를 생성할 수 없습니다", fill='black', font=font)

        return self._process_and_encode_image(img)

    def _process_and_encode_image(self, img):
        """이미지 처리 및 Base64 인코딩"""
        try:
            if img.width > self.max_image_width:
                ratio = self.max_image_width / img.width
                new_height = int(img.height * ratio)
                img = img.resize((self.max_image_width, new_height), Image.LANCZOS)

            buffer = io.BytesIO()
            img.save(buffer, format="JPEG", quality=self.image_quality, optimize=True)
            buffer.seek(0)

            encoded_image = base64.b64encode(buffer.read()).decode('utf-8')
            return f"data:image/jpeg;base64,{encoded_image}"
        except Exception as e:
            logger.error(f"이미지 처리 오류: {e}")
            return ""

    def _process_pptx_alternative(self, pptx_path, folder_name):
        """
        대체 PPTX 처리 방법
        """
        pages_data = []

        for slide_idx in range(5):
            slide_data = {
                "folder_name": folder_name,
                "page_number": str(slide_idx + 1),
                "text_content": f"슬라이드 {slide_idx + 1} - 텍스트 추출 불가",
                "page_summary": f"슬라이드 {slide_idx + 1} - 처리 불가",
                "images": [{
                    'image': self._create_default_image(slide_idx),
                    'description': f'슬라이드 {slide_idx + 1} (자동 생성)',
                    'related_text_ids': []
                }],
                "elements": []
            }

            pages_data.append(slide_data)

        return pages_data